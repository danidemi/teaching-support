#!/usr/bin/env python3
"""Render a deck model (material/slides/session-NN.yml, see
learning-plugin/reference/deck_model_spec.md) into a .pptx, and optionally a .pdf.

This is a lightweight, dependency-minimal renderer: pure Python (python-pptx + PyYAML), no
Docker, no pandoc, no mermaid-cli. Consequences of that choice, stated up front rather than
discovered later:

- A `diagram` body is NOT drawn as a graphic. Its Mermaid source is placed on the slide as text,
  clearly labelled, so nothing is silently lost — draw the actual diagram by hand from that source
  until a graphical Mermaid renderer is wired in.
- A fetched (`source_url`) image is NEVER downloaded by this script. Only a local `asset` file is
  embedded. A `source_url` image is rendered as a labelled placeholder carrying its url, licence,
  attribution and alt text, so a human can fetch and review it deliberately.
- PDF export is best-effort via a locally installed LibreOffice (`soffice`/`libreoffice` on PATH).
  If neither is found, the script still produces the .pptx and says plainly that no PDF was made.

Usage:
    python3 render_deck.py preview <deck.yml> [--pdf]   # always renders, stamps [DRAFT]
    python3 render_deck.py render  <deck.yml> [--pdf]   # refuses unless status: approved
"""

from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)

DRAFT_RED = RGBColor(0xC0, 0x1C, 0x1C)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)


class DeckError(Exception):
    pass


def load_deck(path: Path) -> dict:
    with path.open(encoding="utf-8") as fh:
        deck = yaml.safe_load(fh)
    for key in ("deck", "status", "session", "title", "segments"):
        if key not in deck:
            raise DeckError(f"{path}: missing required top-level key '{key}'")
    if not isinstance(deck["segments"], list) or not deck["segments"]:
        raise DeckError(f"{path}: 'segments' must be a non-empty list")
    return deck


def iter_slides(deck: dict):
    """Yield (segment, slide, index) for every authored slide, in deck order."""
    for segment in deck["segments"]:
        for index, slide in enumerate(segment.get("slides") or [], start=1):
            yield segment, slide, index


def slide_label(segment: dict, slide: dict, index: int) -> str:
    return slide.get("id") or f"{segment.get('title', '?')} slide {index}"


def _blank_slide(prs: Presentation):
    # Layout 6 is "Blank" in the default python-pptx template.
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _set_run(paragraph, text, *, size=18, bold=False, italic=False, color=INK, font="Calibri"):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def _draft_stamp(slide):
    tf = _textbox(slide, SLIDE_W - Inches(2.2), Inches(0.15), Inches(2.0), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _set_run(p, "[DRAFT]", size=14, bold=True, color=DRAFT_RED)


def add_title_slide(prs: Presentation, deck: dict, draft: bool):
    slide = _blank_slide(prs)
    tf = _textbox(slide, MARGIN, Inches(2.6), SLIDE_W - 2 * MARGIN, Inches(1.2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p, deck["title"], size=40, bold=True)

    if deck.get("subtitle"):
        tf2 = _textbox(slide, MARGIN, Inches(3.7), SLIDE_W - 2 * MARGIN, Inches(0.8))
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        _set_run(p2, deck["subtitle"], size=22, italic=True, color=MUTED)

    tf3 = _textbox(slide, MARGIN, Inches(4.6), SLIDE_W - 2 * MARGIN, Inches(0.6))
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    course = deck.get("course", "")
    session = deck.get("session", "")
    _set_run(p3, f"{course} — Sessione {session}", size=16, color=MUTED)

    if draft:
        _draft_stamp(slide)
    return slide


def add_section_slide(prs: Presentation, segment: dict, draft: bool):
    slide = _blank_slide(prs)
    tf = _textbox(slide, MARGIN, Inches(3.2), SLIDE_W - 2 * MARGIN, Inches(1.2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p, segment.get("title", ""), size=32, bold=True)
    if draft:
        _draft_stamp(slide)
    return slide


def _render_body(tf, body: dict):
    kind = body.get("kind", "none")

    if kind == "none":
        return

    if kind == "list":
        for i, item in enumerate(body.get("items") or []):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            prefix = f"{i + 1}. " if body.get("ordered") else "•  "
            _set_run(p, prefix + item, size=22)
        return

    if kind == "code":
        caption = body.get("caption")
        first = True
        if caption:
            p = tf.paragraphs[0]
            _set_run(p, caption, size=16, italic=True, color=MUTED)
            first = False
        for line in (body.get("text") or "").rstrip("\n").split("\n"):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            _set_run(p, line or " ", size=16, font="Consolas")
        return

    if kind == "diagram":
        p = tf.paragraphs[0]
        caption = body.get("caption")
        _set_run(
            p,
            f"[DIAGRAM SOURCE — draw by hand, not rendered by this script]"
            + (f"  {caption}" if caption else ""),
            size=14,
            italic=True,
            color=MUTED,
        )
        for line in (body.get("mermaid") or "").rstrip("\n").split("\n"):
            para = tf.add_paragraph()
            _set_run(para, line or " ", size=14, font="Consolas")
        return

    if kind == "callout":
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p, "“" + (body.get("text") or "") + "”", size=28, italic=True)
        return

    if kind == "placeholder":
        p = tf.paragraphs[0]
        _set_run(p, "PLACEHOLDER", size=20, bold=True, color=DRAFT_RED)
        p2 = tf.add_paragraph()
        _set_run(p2, f"Needs: {body.get('needs', '')}", size=16)
        p3 = tf.add_paragraph()
        _set_run(p3, f"Why: {body.get('why', '')}", size=16, italic=True, color=MUTED)
        return

    if kind == "image":
        asset = body.get("asset")
        p = tf.paragraphs[0]
        if asset:
            # The picture is added directly on the slide by the caller (needs slide, not just
            # tf); here we only add the caption/licence line under it.
            pass
        else:
            _set_run(p, "IMAGE NOT EMBEDDED (fetched image — not downloaded by this script)", size=16, bold=True, color=DRAFT_RED)
            for label, key in (("url", "source_url"), ("license", "license"), ("attribution", "attribution"), ("alt", "alt"), ("reviewed", "reviewed")):
                para = tf.add_paragraph()
                _set_run(para, f"{label}: {body.get(key, '')}", size=14, color=MUTED)
        if body.get("caption"):
            para = tf.add_paragraph()
            _set_run(para, body["caption"], size=14, italic=True, color=MUTED)
        return

    raise DeckError(f"unknown body kind: {kind!r}")


def add_content_slide(prs: Presentation, slide_model: dict, draft: bool):
    slide = _blank_slide(prs)

    head_tf = _textbox(slide, MARGIN, Inches(0.35), SLIDE_W - 2 * MARGIN, Inches(1.1))
    _set_run(head_tf.paragraphs[0], slide_model["headline"], size=26, bold=True)

    body = slide_model.get("body") or {"kind": "none"}
    body_top = Inches(1.6)
    body_h = SLIDE_H - body_top - Inches(0.5)
    body_tf = _textbox(slide, MARGIN, body_top, SLIDE_W - 2 * MARGIN, body_h)

    if body.get("kind") == "image" and body.get("asset"):
        asset_path = Path(body["asset"])
        if asset_path.is_file():
            slide.shapes.add_picture(str(asset_path), MARGIN, body_top, height=Inches(4.2))
        else:
            _set_run(body_tf.paragraphs[0], f"MISSING ASSET FILE: {asset_path}", size=18, bold=True, color=DRAFT_RED)

    _render_body(body_tf, body)

    if slide_model.get("links"):
        link_tf = _textbox(slide, MARGIN, SLIDE_H - Inches(0.5), SLIDE_W - 2 * MARGIN, Inches(0.4))
        text = " · ".join(f"{l.get('label', l.get('url'))}" for l in slide_model["links"])
        _set_run(link_tf.paragraphs[0], text, size=12, color=MUTED)

    if draft:
        _draft_stamp(slide)

    notes = slide_model.get("notes") or {}
    notes_tf = slide.notes_slide.notes_text_frame
    header = f"role: {slide_model.get('role', '?')} | goals: {', '.join(slide_model.get('goals') or [])} | timing_min: {notes.get('timing_min', '?')}"
    notes_tf.text = header
    for field in ("talk", "why_here", "links", "watch_for"):
        p = notes_tf.add_paragraph()
        _set_run(p, f"[{field}] {notes.get(field, '')}", size=12)

    return slide


def build_pptx(deck: dict, draft: bool) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, deck, draft)
    for segment in deck["segments"]:
        add_section_slide(prs, segment, draft)
        for slide_model in segment.get("slides") or []:
            add_content_slide(prs, slide_model, draft)
    return prs


def convert_to_pdf(pptx_path: Path) -> Path | None:
    binary = shutil.which("soffice") or shutil.which("libreoffice")
    if not binary:
        return None
    subprocess.run(
        [binary, "--headless", "--convert-to", "pdf", "--outdir", str(pptx_path.parent), str(pptx_path)],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    pdf_path = pptx_path.with_suffix(".pdf")
    return pdf_path if pdf_path.is_file() else None


def main(argv=None):
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("mode", choices=["preview", "render"])
    parser.add_argument("deck", type=Path)
    parser.add_argument("--pdf", action="store_true", help="also attempt a .pdf via LibreOffice, if installed")
    parser.add_argument("--out-dir", type=Path, default=None, help="defaults to <deck's session dir>/out/")
    args = parser.parse_args(argv)

    try:
        deck = load_deck(args.deck)
    except (DeckError, yaml.YAMLError, FileNotFoundError) as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 1

    is_draft = deck.get("status") != "approved"
    if args.mode == "render" and is_draft:
        print(
            f"error: {args.deck} has status: {deck.get('status')!r}, not 'approved'. "
            "'render' refuses a draft — only a human sets status: approved. Use 'preview' instead.",
            file=sys.stderr,
        )
        return 1

    out_dir = args.out_dir or (args.deck.parent / "out")
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = out_dir / (args.deck.stem + ".pptx")

    prs = build_pptx(deck, draft=is_draft)
    prs.save(pptx_path)
    print(f"wrote {pptx_path}")

    if args.pdf:
        pdf_path = convert_to_pdf(pptx_path)
        if pdf_path:
            print(f"wrote {pdf_path}")
        else:
            print(
                "no PDF written: neither 'soffice' nor 'libreoffice' found on PATH. "
                "Install LibreOffice, or open the .pptx and export to PDF by hand.",
                file=sys.stderr,
            )

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
